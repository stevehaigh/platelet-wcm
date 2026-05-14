"""
Build the lab-meeting deck directly with python-pptx, using
~/Desktop/steve-academic-theme.pptx as the template (it has custom layouts
named DEFAULT / TITLE / SECTION / CONTENT / TUFTE / CALLOUT / EQUATION /
TWOCOL / BLANK that pandoc can't address through Quarto's pptx pipeline).

Eight slides, informal lab-meeting tone, same content arc as the Mike report.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

TEMPLATE = Path.home() / "Desktop" / "steve-academic-theme.pptx"
OUT      = Path(__file__).parent / "lab-meeting-2026-05-14.pptx"
FIGURES  = Path(__file__).parent.parent / "figures"

# Layout indices (from inspecting the template)
L_TITLE   = 1
L_CONTENT = 3
L_CALLOUT = 5
L_TWOCOL  = 7
L_BLANK   = 8


def set_ph(slide, idx, text):
	"""Set placeholder text by placeholder idx."""
	for ph in slide.placeholders:
		if ph.placeholder_format.idx == idx:
			ph.text = text
			return
	raise KeyError(f"placeholder idx={idx} not found on slide layout {slide.slide_layout.name!r}")


def set_ph_bullets(slide, idx, bullets):
	"""Set placeholder with one bullet per item (each item is a string)."""
	for ph in slide.placeholders:
		if ph.placeholder_format.idx == idx:
			tf = ph.text_frame
			tf.text = bullets[0]
			for line in bullets[1:]:
				p = tf.add_paragraph()
				p.text = line
			return
	raise KeyError(f"placeholder idx={idx} not found")


prs = Presentation(str(TEMPLATE))

# Template ships with 8 demo slides showing each layout. Remove them before
# we add our own; layout masters stay attached. Drop relationships + parts
# too, otherwise the saved .pptx ends up with duplicate zip entries.
from pptx.oxml.ns import qn
xml_slides = prs.slides._sldIdLst
slides_list = list(xml_slides)
for sld in slides_list:
	rId = sld.get(qn("r:id"))
	prs.part.drop_rel(rId)
	xml_slides.remove(sld)

# ── Slide 1: Title ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[L_TITLE])
set_ph(s, 102, "Building a platelet whole-cell model")
set_ph(s, 103, "What worked, what didn't, and the k₃ story\nSteve Haigh  ·  lab meeting  ·  May 2026")

# ── Slide 2: Why a whole-cell framework? ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[L_TWOCOL])
set_ph(s, 101, "Framework choice")
set_ph(s, 102, "Three options. wcEcoli won.")
set_ph(s, 103, "What I evaluated")
set_ph_bullets(s, 104, [
	"COPASI — great ODE solver, SBML import. But single-pathway focus; no copy-number accounting; no mass balance across compartments.",
	"SBML / BioModels — a catalogue and a format, not an engine. Models in isolation; re-use is integration cost.",
	"wcEcoli (Covert Lab, 2020) — the only validated multi-process whole-cell model. Mass-balanced. Every protein counted. Compartmentalised. But: E. coli-specific, ~80 kloc, unmaintained since 2022.",
])
set_ph(s, 105, "What I did")
set_ph_bullets(s, 106, [
	"Forked wcEcoli.",
	"Pruned all E. coli biology (~1 month).",
	"Rebuilt reconstruction and processes from the platelet literature.",
	"Kept the simulation engine, state partitioning, listener / analysis framework — ~5,000 lines of validated infrastructure I didn't have to write.",
])

# ── Slide 3: The process ───────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
set_ph(s, 101, "Methodology")
set_ph(s, 102, "How a biologist builds a calibrated model")
set_ph_bullets(s, 103, [
	"For each of 12 mechanisms, a fixed five-step process:",
	"1. Anchor paper. Dolan & Diamond 2014 supplied the validation experiment (Fig. 4 Ca²⁺ transients ± extracellular Ca²⁺) and resting-state targets (100 nM cyt, 250 µM DTS).",
	"2. Literature review for kinetics. Primary sources only: deYoung-Keizer 1992 for IP3R, Caride 2007 for PMCA, Hoover & Lewis 2011 for SOCE, Mazet 2020 for the PI cycle...",
	"3. Species enumeration. Every Ca²⁺-binding or -gating protein state, with a compartment tag. ~50 species across the calcium pathway.",
	"4. Copy numbers. Burkhart 2012 platelet proteomics.",
	"5. Rate constants. Primary-source values, units normalised, every value carrying a citation comment in code.",
	"Then: does the resting state hold? Do the timescales line up? Do the unit tests still pass?",
])

# ── Slide 4: The Purvis k3 story (CALLOUT) ─────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[L_CALLOUT])
set_ph(s, 101, "Methodology / sanity-checks")
set_ph(s, 102, "When the sanity-checks earn their keep")
set_ph(s, 104, "The Purvis 2008 k₃ story")
set_ph_bullets(s, 105, [
	"After implementing SERCA from Purvis 2008 Table 1: resting state diverged — cytosolic Ca²⁺ ran away or collapsed depending on initial conditions.",
	"Back-traced the flux balance to one rate constant: k₃, the E₁P·Ca → E₂P·Ca phosphoryl-transfer step.",
	"Cross-checked Purvis 2008 Table 1 against the primary source it cited (Dode 2002). The Purvis value was the reciprocal of the Dode value — a transcription error propagated unchecked.",
	"Corrected the value; the resting state held. Three other Purvis values spot-checked against primaries; one more rounding-induced drift found.",
	"Moral: every numerical value carries a primary-source citation in the code. AI would have happily reproduced the typo.",
])

# ── Slide 5: Model diagram (BLANK + manual title + image) ──────────────────
s = prs.slides.add_slide(prs.slide_layouts[L_BLANK])
# Title box (matches CONTENT layout's title placement)
tb = s.shapes.add_textbox(Inches(0.75), Inches(0.30), Inches(11.83), Inches(0.30))
tb.text_frame.text = "Model"
tb.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
tb = s.shapes.add_textbox(Inches(0.75), Inches(0.6), Inches(11.83), Inches(0.65))
tb.text_frame.text = "Twelve coupled mechanisms, agonist → Ca²⁺ peak"
tb.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
tb.text_frame.paragraphs[0].runs[0].font.bold = True
# Diagram
img = FIGURES / "model-status-graphviz.png"
s.shapes.add_picture(str(img), Inches(0.5), Inches(1.5), height=Inches(5.7))

# ── Slide 6: Validation (TWOCOL) ───────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[L_TWOCOL])
set_ph(s, 101, "Result")
set_ph(s, 102, "Validation: Dolan & Diamond 2014 Fig 4")
set_ph(s, 103, "Acceptance criteria")
set_ph_bullets(s, 104, [
	"Resting cytosolic Ca²⁺ within 100 ± 10 nM band",
	"Resting DTS Ca²⁺ in the literature range",
	"IP₃-stimulated transient peak height in band",
	"Transient duration matches Fig 4 shape",
	"Paired ± extracellular Ca²⁺ condition difference reproduces",
])
set_ph(s, 105, "Model state at v0.4.1")
set_ph_bullets(s, 106, [
	"Resting cyt Ca²⁺: 104 nM  ✓",
	"Resting DTS Ca²⁺: 235 µM  ✓",
	"Resting IP₃: 50 nM  ✓",
	"Resting Gαq-active: 100 of 5000  ✓",
	"5 / 5 Phase 3 acceptance criteria",
	"21 / 21 unit tests pass",
	"Driven by physiological agonists — 1 nM thrombin, 10 µM ADP. No hand-fitted IP₃ forcing.",
])

# ── Slide 7: Result figure (BLANK + manual title + image) ──────────────────
s = prs.slides.add_slide(prs.slide_layouts[L_BLANK])
tb = s.shapes.add_textbox(Inches(0.75), Inches(0.30), Inches(11.83), Inches(0.30))
tb.text_frame.text = "Result"
tb.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
tb = s.shapes.add_textbox(Inches(0.75), Inches(0.6), Inches(11.83), Inches(0.65))
tb.text_frame.text = "Free vs bound Ca²⁺ during an IP₃-driven transient"
tb.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
tb.text_frame.paragraphs[0].runs[0].font.bold = True
img = FIGURES / "ca-bound-free-v0.4.0.png"
s.shapes.add_picture(str(img), Inches(1.0), Inches(1.5), height=Inches(5.5))
# Brief caption under image
tb = s.shapes.add_textbox(Inches(0.75), Inches(7.0), Inches(11.83), Inches(0.4))
tb.text_frame.text = "Cyt rises ~100 nM → 500 nM over 60 s; DTS depletes ~60 % and refills via SOCE; buffer occupancy tracks the free pools."
tb.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
tb.text_frame.paragraphs[0].runs[0].font.italic = True

# ── Slide 8: Where it goes from here ───────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
set_ph(s, 101, "Outlook")
set_ph(s, 102, "Where it goes from here")
set_ph_bullets(s, 103, [
	"Near-term biology:",
	"  Granule release — the model currently stops at the Ca²⁺ peak; v0.5 lets that peak do something (dense and α-granule exocytosis).",
	"  Long-recovery cytosolic Ca²⁺ collapse — a deYoung-Keizer bistability artefact at sustained stim > 2000 s, fixable.",
	"  P2Y₁₂ / Gi pathway — the inhibitory ADP arm (clopidogrel target), parallel to Gq.",
	"",
	"Cross-disciplinary:",
	"  Platelet – tumour interactions in metastasis. Calibrated platelet model could in principle be coupled to a tumour-cell model — relevant to this lab's work.",
	"  Methodology generalises to other single-cell calibration problems.",
])

prs.save(str(OUT))
print(f"Wrote {OUT} — {len(prs.slides)} slides")
